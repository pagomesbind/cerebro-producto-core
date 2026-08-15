# -*- coding: utf-8 -*-
"""Rough approximate renderer for python-pptx decks, for use when no
LibreOffice is available (see SKILL.md Paso 6 — this is what to reach for
instead of soffice.py in that situation). NOT pixel-accurate — no real OOXML
layout engine, no kerning, no real word-wrap metrics — just enough to catch
gross layout bugs (overlaps, off-canvas shapes, wrong colors, missing
assets) before a slide reaches the user. Always still ask the user to open
the real .pptx in PowerPoint for the actual visual sign-off.

Handles: rects/roundrects/ellipses with solid fill, straight connectors
(lines), pictures (alpha-composited), and text (word-wrapped per paragraph
within the shape's width, using the real brand fonts so sizes are at least
in the right ballpark).

Usage: python approx_render.py <deck.pptx> <out_dir> [file_prefix]
Then tile the resulting PNGs into a contact sheet and view it — see
SKILL.md Paso 6 for the exact snippet.
"""
import sys, os
import io
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

FONT_DIR = os.path.join(os.path.dirname(__file__), "..", "assets", "fonts")
FONT_FILES = {
    "Figtree": "Figtree-Regular.ttf",
    "Figtree Bold": "Figtree-Bold.ttf",
    "Figtree Light": "Figtree-Light.ttf",
    "Figtree Medium": "Figtree-Medium.ttf",
    "Inter": "Inter-Regular.ttf",
}
_font_cache = {}
def get_font(name, bold, size_px):
    key = (name, bold, size_px)
    if key in _font_cache:
        return _font_cache[key]
    fname = FONT_FILES.get(name, "Figtree-Regular.ttf")
    if bold and name == "Figtree":
        fname = "Figtree-Bold.ttf"
    if bold and name == "Inter":
        fname = "Inter-Bold.ttf"
    path = os.path.join(FONT_DIR, fname)
    try:
        f = ImageFont.truetype(path, size_px)
    except Exception:
        f = ImageFont.load_default()
    _font_cache[key] = f
    return f

SCALE = 130  # px per inch

def emu_to_px(v):
    return int(Emu(v) / 914400 * SCALE)

def render(pptx_path, out_dir, prefix="slide"):
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    W = emu_to_px(prs.slide_width)
    H = emu_to_px(prs.slide_height)
    paths = []
    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), (255, 255, 255))
        # background
        try:
            bg = slide.background.fill
            if bg.type == 1:
                img.paste(tuple(bg.fore_color.rgb) if False else _rgb(bg.fore_color.rgb), (0, 0, W, H))
        except Exception:
            pass
        draw = ImageDraw.Draw(img, "RGBA")
        for shp in slide.shapes:
            _draw_shape(img, draw, shp)
        p = os.path.join(out_dir, f"{prefix}_{idx:02d}.png")
        img.save(p)
        paths.append(p)
    return paths

def _rgb(rgbcolor):
    return (rgbcolor[0], rgbcolor[1], rgbcolor[2])

def _draw_shape(img, draw, shp):
    x, y = emu_to_px(shp.left or 0), emu_to_px(shp.top or 0)
    w, h = emu_to_px(shp.width or 0), emu_to_px(shp.height or 0)
    try:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic = Image.open(io.BytesIO(shp.image.blob)).convert("RGBA")
            pic = pic.resize((max(w,1), max(h,1)))
            img.paste(pic, (x, y), pic)
            return
    except Exception:
        draw.rectangle([x, y, x + w, y + h], outline=(255, 0, 0))
    # connector/line
    if str(type(shp).__name__) == "Connector":
        try:
            col = _rgb(shp.line.color.rgb)
        except Exception:
            col = (150, 150, 150)
        draw.line([x, y, x + w, y + h], fill=col, width=2)
        return
    # autoshape / freeform: fill
    try:
        if shp.fill.type == 1:
            col = _rgb(shp.fill.fore_color.rgb)
            if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shp.auto_shape_type is not None and "OVAL" in str(shp.auto_shape_type):
                draw.ellipse([x, y, x + w, y + h], fill=col + (255,))
            else:
                draw.rounded_rectangle([x, y, x + w, y + h], radius=min(w, h) * 0.08, fill=col + (255,))
    except Exception:
        pass
    # text — word-wrapped within the shape width (crude but honest about overflow)
    if shp.has_text_frame:
        ty = y + 4
        maxw = max(w - 8, 10)
        for para in shp.text_frame.paragraphs:
            align = str(para.alignment) if para.alignment else "LEFT"
            words = []  # list of (word_text, font, color), with explicit space tokens
            for r in para.runs:
                if not r.text:
                    continue
                sz = r.font.size.pt if r.font.size else 12
                px = max(int(sz * SCALE / 72 * 0.8), 8)
                fnt = get_font(r.font.name or "Figtree", r.font.bold, px)
                col = _rgb(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else (0,0,0)
                for tok in r.text.split(" "):
                    if tok:
                        words.append((tok, fnt, col))
                    words.append((" ", fnt, col))
            if not words:
                ty += 14
                continue
            # greedy wrap into lines
            lines = [[]]
            cur_w = 0
            for tok, fnt, col in words:
                bbox = draw.textbbox((0, 0), tok, font=fnt)
                tw = bbox[2] - bbox[0]
                if cur_w + tw > maxw and lines[-1]:
                    lines.append([])
                    cur_w = 0
                lines[-1].append((tok, fnt, col, tw))
                cur_w += tw
            for line in lines:
                if not line:
                    ty += 14
                    continue
                line_h = max(draw.textbbox((0, 0), t, font=f)[3] for t, f, c, tw in line if t.strip()) if any(t.strip() for t,f,c,tw in line) else 12
                total_w = sum(tw for t, f, c, tw in line)
                if "CENTER" in align:
                    cx = x + max((w - total_w) // 2, 0)
                elif "RIGHT" in align:
                    cx = x + max(w - total_w - 4, 0)
                else:
                    cx = x + 4
                for tok, fnt, col, tw in line:
                    draw.text((cx, ty), tok, font=fnt, fill=col)
                    cx += tw
                ty += line_h + 5

if __name__ == "__main__":
    render(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else "slide")
    print("done")
