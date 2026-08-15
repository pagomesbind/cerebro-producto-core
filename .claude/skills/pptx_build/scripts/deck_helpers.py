"""
Reusable building blocks for Bind PSP-branded decks (python-pptx).

This is the **house style** — reverse-engineered from the user's own final,
hand-finished deck (`raw/onboarding_estrategico_reunion_entendimiento.pptx`,
confirmed 2026-08-10 as the reference for all future decks), not from
`template_bind.pptx` (marketing's template) directly. The template was the
starting point for colors/fonts/logo/flower assets, but the user judged its
own slide layouts too rigid/marketing-oriented for internal work — this
module encodes what they actually kept, changed, and simplified after using
it once. Read references/design-system.md before assuming anything here is
still template-accurate; several defaults below deliberately diverge from
the template (no page-number badge, no rule under the title, a fixed
"BIND PSP | <deck name>" kicker instead of a per-section one).

This skill is scoped to PPTX only — there is no PDF/other-medium sibling
system to keep in sync with here anymore.

Usage sketch:

    from deck_helpers import *

    prs = new_presentation()
    set_deck_name("Onboarding Estratégico")   # drives the kicker on every content slide

    cover_slide(prs, eyebrow="Onboarding Estratégico", title="Reunión de entendimiento")

    section_break_slide(prs, "BLOQUE 1", "Situación actual y decisión estratégica",
                         variant="dark")

    s = light_slide(prs)
    page_title(s, "El problema", icon="ico-05")
    card_row(s, Inches(0.67), Inches(1.72), Inches(8.62), Inches(1.43), [
        dict(title="Mandato del banco", body="Texto del cuerpo..."),
        dict(title="El punto de partida", body="..."),
        dict(title="Por qué también nos conviene", body="..."),
    ])
    ...
    prs.save("outputs/mi_deck.pptx")
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- brand tokens (Bind PSP PPTX system — see references/brand-bind-psp.md) ----------
AMARILLO = RGBColor(0xF5, 0xF4, 0x00)
AMARILLO_TINT = RGBColor(0xFD, 0xFD, 0xC7)     # "FDFDC7" pale-yellow card fill
AZUL = RGBColor(0x4B, 0x64, 0xE6)              # the ONLY blue — normalize away any 4285F4 drift
AZUL_DARK = RGBColor(0x33, 0x27, 0x7F)
NEGRO = RGBColor(0x00, 0x00, 0x00)
GRIS = RGBColor(0x99, 0x99, 0x99)          # dividers, muted/secondary text
SURFACE = RGBColor(0xF3, 0xF3, 0xF3)
GRIS_PANEL = RGBColor(0xEF, 0xEF, 0xEF)    # closing-slide panel gray
PANEL_OSCURO = RGBColor(0x22, 0x22, 0x22)  # dark title/section-slide background
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

# Real brand fonts (see references/brand-bind-psp.md "Fonts"). "Figtree Light"/
# "Figtree Medium" are their own registered family names, not a weight flag on
# "Figtree" — pass them verbatim to font.name, never bold=True to fake them.
FONT_HEAD = "Figtree"
FONT_HEAD_LIGHT = "Figtree Light"
FONT_HEAD_MEDIUM = "Figtree Medium"        # secondary/nested card headings (see card() title_font)
FONT_BODY = "Inter"              # general-purpose body copy — rarely needed, see FONT_CARD_BODY
FONT_CARD_BODY = "Figtree Light"  # card/paragraph body font — use this for anything meant to
                                   # look native to this deck, not Inter.

import os
_ASSETS = os.path.join(os.path.dirname(__file__), "..", "assets")
LOGO_LIGHT = os.path.join(_ASSETS, "logo-wordmark-light.png")   # black text — light backgrounds
LOGO_DARK = os.path.join(_ASSETS, "logo-wordmark-dark.png")     # white text — dark (#222222) bg
LOGO_BLACK = os.path.join(_ASSETS, "logo-wordmark-black.png")   # monochrome black — on yellow bg
FLOWER_YELLOW = os.path.join(_ASSETS, "flower-yellow.png")
FLOWER_DARK = os.path.join(_ASSETS, "flower-dark.png")          # black flower — on light bg
FLOWER_WHITE = os.path.join(_ASSETS, "flower-white.png")        # white flower — on yellow bg
ICONS_DIR = os.path.join(_ASSETS, "icons")

_FLOWER_BY_NAME = {"yellow": FLOWER_YELLOW, "dark": FLOWER_DARK, "white": FLOWER_WHITE}
_LOGO_BY_VARIANT = {"light": LOGO_LIGHT, "dark": LOGO_DARK, "black": LOGO_BLACK}

# The deck's display name, used by every content-slide kicker (see set_deck_name()).
_DECK_NAME = [None]


def set_deck_name(name):
    """Call once per deck, before building content slides. Every content
    slide's kicker reads 'BIND PSP | <name>' — confirmed 2026-08-10 as the
    fixed, deck-wide label (replaces the old per-section kicker)."""
    _DECK_NAME[0] = name


def icon_path(name, variant="black"):
    """name like 'ico-05' (see assets/icons/INDEX.md + INDEX.png to pick one).
    variant is 'black' or 'white' — pick whichever contrasts with where you're
    placing it. icon_badge()/card() below usually resolve this for you."""
    return os.path.join(ICONS_DIR, variant, f"{name}.png")


def _is_light(rgb):
    """Perceived luminance check — used to auto-pick a black/white icon or
    text color against a given fill so it never disappears into its own
    background."""
    r, g, b = rgb[0], rgb[1], rgb[2]
    return (0.299 * r + 0.587 * g + 0.114 * b) > 150


def new_presentation():
    """10 x 5.625in (16:9) canvas — every coordinate in this module is
    calibrated to it. See design-system.md 'Canvas'."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


# ---------------------------------------------------------------------------
# Text primitives
# ---------------------------------------------------------------------------

def textbox(slide, x, y, w, h, text, size, color, bold=False, font=FONT_BODY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """Single run, single paragraph — for plain labels. For two-tone titles
    (bold run + light run) use multi_text()."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return tb


def multi_text(slide, x, y, w, h, paragraphs, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_spacing=1.05, space_after=0):
    """One textbox holding several paragraphs, each made of several runs with
    independent font/size/weight/color — e.g. a title's bold lead-in run
    followed by a lighter continuation run, or a heading paragraph followed
    by a lighter subtitle paragraph.

    paragraphs: list of paragraphs; each paragraph is a list of run-dicts:
        {"text": str, "font": FONT_HEAD, "bold": False, "size": 40, "color": NEGRO}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(space_after)
        for rd in runs:
            r = p.add_run()
            r.text = rd["text"]
            r.font.size = Pt(rd.get("size", 18))
            r.font.bold = rd.get("bold", False)
            r.font.name = rd.get("font", FONT_HEAD)
            r.font.color.rgb = rd.get("color", NEGRO)
    return tb


def bullets(slide, x, y, w, h, items, size=13, color=NEGRO, font=FONT_CARD_BODY,
            line_spacing=1.25, space_after=10, bullet_char="•  "):
    """Plain bulleted paragraph list. items: list of str, or list of
    (bold_lead, rest) tuples. Prefer list_rows() instead when the content is
    a short enumerated list of risks/steps/decisions — it reads much better
    (see design-system.md) — reach for this only for genuine prose bullets
    that don't fit that shape."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = bullet_char + lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r1.font.name = FONT_HEAD
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.bold = False
            r2.font.color.rgb = color
            r2.font.name = font
        else:
            r = p.add_run()
            r.text = bullet_char + item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = font
    return tb


def list_rows(slide, x, y, w, items, marker_color=AZUL, lead_color=None, body_color=NEGRO,
              row_h=Inches(0.7), gap=Inches(0.08), marker_d=Inches(0.394), size=14):
    """Enumerated list pattern for risks / next-steps / decisions: a small
    solid circle marker + a bold colored lead phrase + a light body,
    left-aligned in a row, one row per item. This is the confirmed
    replacement for plain bullets() on that kind of content — much more
    legible at a glance. items: list of (lead, rest) tuples.
    Returns the y (Emu) just after the last row, for stacking more content."""
    if lead_color is None:
        lead_color = marker_color
    cy = int(y)
    marker_x = int(x)
    text_x = marker_x + int(marker_d) + int(Inches(0.18))
    text_w = int(x) + int(w) - text_x
    for lead, rest in items:
        marker_y = cy + (int(row_h) - int(marker_d)) // 2
        icon_badge(slide, Emu(marker_x), Emu(marker_y), marker_d, bg=marker_color)
        tb = slide.shapes.add_textbox(Emu(text_x), Emu(cy), Emu(text_w), row_h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        r1 = p.add_run()
        r1.text = lead
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = lead_color
        r1.font.name = FONT_HEAD
        r2 = p.add_run()
        r2.text = rest
        r2.font.size = Pt(size)
        r2.font.bold = False
        r2.font.color.rgb = body_color
        r2.font.name = FONT_CARD_BODY
        cy += int(row_h) + int(gap)
    return Emu(cy)


# ---------------------------------------------------------------------------
# Logo / flower (real brand assets, not drawn shapes)
# ---------------------------------------------------------------------------

def add_logo(slide, variant="light", x=Inches(1.29), y=Inches(3.31), height=Inches(0.85)):
    """Big/free-standing logo placement (cover, section-break panel). For the
    small recurring corner mark on content slides use add_logo_small()."""
    return slide.shapes.add_picture(_LOGO_BY_VARIANT[variant], x, y, height=height)


def add_logo_small(slide, variant="light"):
    """The small logo mark every content slide carries, top-left, at a fixed
    position — part of add_header()."""
    return slide.shapes.add_picture(_LOGO_BY_VARIANT[variant], Inches(0.17), Inches(0.13),
                                     width=Inches(0.86), height=Inches(0.32))


def add_flower(slide, variant="yellow", x=Inches(7.91), y=Inches(0.64), width=Inches(2.09),
                height=Inches(4.18), crop_left=0.499, crop_right=0.001):
    """Decorative flower mark. variant must contrast with the panel it sits
    on: dark panel -> yellow flower, yellow panel -> white flower, light/gray
    panel -> yellow flower. section_break_slide() applies this automatically.

    Default geometry (2.09x4.18, cropped to the right half of the source
    square PNG) is deliberate — the full square mark is wide enough to run
    straight into a title block placed anywhere near it. Don't widen this or
    drop the crop without checking clearance against nearby text."""
    pic = slide.shapes.add_picture(_FLOWER_BY_NAME[variant], x, y, width=width, height=height)
    pic.crop_left = crop_left
    pic.crop_right = crop_right
    return pic


def add_kicker(slide, dark_bg=False):
    """Top-right running label — fixed 'BIND PSP | <deck name>' (set via
    set_deck_name()), not a per-section label. Two runs: bold 'BIND PSP '
    + light '| <NAME>'."""
    color = BLANCO if dark_bg else GRIS
    name = (_DECK_NAME[0] or "").upper()
    multi_text(slide, Inches(4.6), Inches(0.2), Inches(5.13), Inches(0.28),
               [[{"text": "BIND PSP ", "font": FONT_HEAD, "bold": True, "size": 9, "color": color},
                 {"text": f"| {name}", "font": FONT_HEAD_LIGHT, "bold": False, "size": 9, "color": color}]],
               align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_header(slide, dark_bg=False, logo_variant=None):
    """Standard content-slide header: small logo top-left + kicker top-right.
    No page-number badge — that element was dropped (see design-system.md).
    page_title() already calls this; use it standalone only if a slide skips
    page_title() for a custom title treatment."""
    if logo_variant is None:
        logo_variant = "dark" if dark_bg else "light"
    add_logo_small(slide, variant=logo_variant)
    add_kicker(slide, dark_bg=dark_bg)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def dark_slide(prs):
    """Bare full #222222 panel, no header — use for the cover or a custom
    dark composition. For a titled section break use section_break_slide()."""
    s = _blank_slide(prs)
    set_background(s, PANEL_OSCURO)
    return s


def light_slide(prs):
    """Content slide — plain white background, no header. Call page_title()
    right after (it applies the header + title together) unless the slide
    needs a fully custom top treatment."""
    s = _blank_slide(prs)
    set_background(s, BLANCO)
    return s


def cover_slide(prs, eyebrow, title, subtitle=None, logo_variant="dark"):
    """Deck opener — dark panel, flower inset right, eyebrow+title in two
    runs (bold amarillo + light white), rule, logo."""
    s = dark_slide(prs)
    add_flower(s, variant="yellow", x=Inches(7.91), y=Inches(0.64), height=Inches(4.18))
    paras = [[{"text": eyebrow, "font": FONT_HEAD, "bold": True, "size": 40, "color": AMARILLO}]]
    if title:
        paras.append([{"text": title, "font": FONT_HEAD_LIGHT, "bold": False, "size": 26, "color": BLANCO}])
    multi_text(s, Inches(1.15), Inches(1.53), Inches(6.3), Inches(1.7), paras, line_spacing=1.08)
    rule_y = Inches(3.2)
    divider(s, Inches(1.29), rule_y, Inches(3.24), color=AMARILLO, weight=1.0)
    if subtitle:
        textbox(s, Inches(1.29), Inches(3.32), Inches(5.5), Inches(0.6), subtitle, 12, GRIS,
                font=FONT_CARD_BODY, line_spacing=1.25)
    add_logo(s, variant=logo_variant, x=Inches(1.29), y=Inches(3.85) if subtitle else Inches(3.31),
             height=Inches(0.6))
    return s


_SECTION_VARIANTS = {
    # bg, flower, title_color, subtitle_color, divider_color, logo_variant
    "dark": (PANEL_OSCURO, "yellow", AMARILLO, BLANCO, BLANCO, "dark"),
    "yellow": (AMARILLO, "white", NEGRO, NEGRO, NEGRO, "black"),
    "gray": (SURFACE, "yellow", NEGRO, NEGRO, NEGRO, "light"),
}


def section_break_slide(prs, block_label, title, subtitle=None, variant="dark"):
    """Quiebre de sección / block divider — inset panel + flower + eyebrow
    (block label) + title + rule + small logo. 3 fill variants (dark/yellow/
    gray), each with its own contrast-correct flower/text/logo combo — don't
    mix them across a variant, see design-system.md."""
    bg, flower, title_c, sub_c, div_c, logo_v = _SECTION_VARIANTS[variant]
    s = _blank_slide(prs)
    set_background(s, BLANCO)
    panel = rounded_card(s, Inches(0.22), Inches(0.24), Inches(9.56), Inches(5.14), bg)
    panel.adjustments[0] = 0.02
    add_flower(s, variant=flower, x=Inches(0.22), y=Inches(0.72))
    tx = Inches(3.05)
    paras = [[{"text": block_label, "font": FONT_HEAD, "bold": True, "size": 15, "color": title_c}],
             [{"text": title, "font": FONT_HEAD, "bold": True, "size": 32, "color": title_c}]]
    if subtitle:
        paras.append([{"text": subtitle, "font": FONT_HEAD_LIGHT, "bold": False, "size": 15, "color": sub_c}])
    multi_text(s, tx, Inches(1.75), Inches(6.3), Inches(2.2), paras, line_spacing=1.08, space_after=6)
    divider(s, tx, Inches(3.55), Inches(3.2), color=div_c, weight=1.0)
    add_logo(s, variant=logo_v, x=tx, y=Inches(3.72), height=Inches(0.42))
    return s


def closing_slide(prs, message="¡Gracias!", subtitle=None):
    """Closing slide — light gray panel, big Figtree-bold message, logo/
    flower accent."""
    s = _blank_slide(prs)
    set_background(s, BLANCO)
    panel = rounded_card(s, Inches(0.45), Inches(0.47), Inches(9.12), Inches(4.76), GRIS_PANEL)
    panel.adjustments[0] = 0.03
    textbox(s, Inches(1.32), Inches(2.05), Inches(5.5), Inches(1.4), message, 54, NEGRO,
            bold=True, font=FONT_HEAD)
    if subtitle:
        textbox(s, Inches(1.32), Inches(3.05), Inches(5.5), Inches(0.5), subtitle, 14, GRIS,
                font=FONT_CARD_BODY)
    add_flower(s, variant="yellow", x=Inches(7.48), y=Inches(0.47), height=Inches(2.08))
    add_logo(s, variant="light", x=Inches(6.2), y=Inches(2.5), height=Inches(0.55))
    return s


# ---------------------------------------------------------------------------
# Content components
# ---------------------------------------------------------------------------

def page_title(slide, text, icon, sub=None, dark_bg=False):
    """Every content slide's heading: header (logo+kicker) + a required
    F5F400 icon-badge to the left of a 43pt Figtree Bold title. No rule
    under the title (confirmed dropped — more breathing room before content
    starts instead). `icon` is required now: every content slide carries a
    topic-relevant icon (see assets/icons/INDEX.md) — pick one that fits,
    don't leave the slide without one."""
    add_header(slide, dark_bg=dark_bg)
    icon_badge(slide, Inches(0.669), Inches(0.793), Inches(0.615), bg=AMARILLO, icon=icon,
               icon_scale=0.64)
    title_color = BLANCO if dark_bg else NEGRO
    textbox(slide, Inches(1.378), Inches(0.62), Inches(6.93), Inches(0.93), text, 43, title_color,
            bold=True, font=FONT_HEAD)
    if sub:
        textbox(slide, Inches(1.378), Inches(1.5), Inches(7.5), Inches(0.4), sub, 12, GRIS,
                font=FONT_CARD_BODY)


def divider(slide, x, y, w, color=GRIS, weight=1.0):
    """Thin horizontal rule — still used between grid columns / under the
    cover and section-break titles, just not under the page title anymore."""
    ln = slide.shapes.add_connector(1, x, y, Emu(int(x) + int(w)), y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def rounded_card(slide, x, y, w, h, fill_color):
    """No drop shadow. Corner radius ~0.056 of the shorter side."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.056
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    no_line(shp)
    return shp


def icon_badge(slide, cx, cy, d, bg=AZUL, icon=None, icon_variant=None, icon_scale=0.55,
                label=None, label_color=None, size=16):
    """Colored circle marker. Pass icon='ico-NN' (see assets/icons/INDEX.md)
    for a real brand icon — variant (black/white) is auto-picked for
    contrast against `bg` unless overridden. Pass a short label instead only
    as a placeholder while iterating, or omit both for a plain solid dot
    (list_rows() marker)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    no_line(shp)
    if icon:
        variant = icon_variant or ("black" if _is_light(bg) else "white")
        isz = Emu(int(int(d) * icon_scale))
        ix = Emu(int(cx) + (int(d) - int(isz)) // 2)
        iy = Emu(int(cy) + (int(d) - int(isz)) // 2)
        slide.shapes.add_picture(icon_path(icon, variant), ix, iy, width=isz, height=isz)
    elif label:
        if label_color is None:
            label_color = NEGRO if _is_light(bg) else BLANCO
        tf = shp.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = label_color
        r.font.name = FONT_HEAD
    return shp


# Rotation used by card_row() when no explicit fills are given — the
# confirmed "3 things, 3 different colors" pattern, not a uniform gray grid.
CARD_PALETTE = [AZUL, AMARILLO_TINT, AMARILLO, SURFACE]


def card(slide, x, y, w, h, fill=SURFACE, icon=None, icon_bg=None, title=None,
         title_color=None, title_size=14, title_font=FONT_HEAD, body=None,
         body_color=None, body_size=9, pad=Inches(0.2)):
    """The base content card: rounded rect, an icon-in-circle badge inset
    top-left (SAME color as the card fill — it's meant to blend in, not
    contrast — the icon glyph itself is what contrasts), a bold title, a
    light body. Reach for card_row() instead of calling this directly when
    laying out 2-4 cards side by side (handles width/gap/palette for you)."""
    box = rounded_card(slide, x, y, w, h, fill)
    if title_color is None:
        title_color = BLANCO if fill in (AZUL, PANEL_OSCURO) else NEGRO
    if body_color is None:
        body_color = BLANCO if fill in (AZUL, PANEL_OSCURO) else GRIS
    inner_x = Emu(int(x) + int(pad))
    inner_w = Emu(int(w) - 2 * int(pad))
    cursor_y = int(y) + int(pad)
    if icon:
        d = Inches(0.5)
        icon_badge(slide, inner_x, Emu(cursor_y), d, bg=icon_bg or fill, icon=icon)
        cursor_y += int(d) + int(Inches(0.1))
    if title:
        th = Inches(0.4)
        textbox(slide, inner_x, Emu(cursor_y), inner_w, th, title, title_size, title_color,
                bold=True, font=title_font, line_spacing=1.05)
        cursor_y += int(th)
    if body:
        bh = int(y) + int(h) - int(pad) - cursor_y
        textbox(slide, inner_x, Emu(cursor_y), inner_w, Emu(max(bh, 0)), body, body_size,
                body_color, font=FONT_CARD_BODY, line_spacing=1.28)
    return box


def card_row(slide, x, y, w_total, h, items, gap=Inches(0.2), palette=None, icon_default=None):
    """Lay out N cards left-to-right, auto-computing each card's width and
    auto-coloring from `palette` (defaults to CARD_PALETTE) unless an item
    sets its own fill. This is the "explain 3 things" / "explain N parallel
    options" pattern — confirmed: vary the color across the row, don't make
    them all the same gray.

    items: list of dicts — {title, body, icon (optional), fill (optional),
    title_color/body_color (optional overrides)}."""
    palette = palette or CARD_PALETTE
    n = len(items)
    w = Emu((int(w_total) - (n - 1) * int(gap)) // n)
    cx = int(x)
    for i, item in enumerate(items):
        fill = item.get("fill", palette[i % len(palette)])
        card(slide, Emu(cx), y, w, h, fill=fill, icon=item.get("icon", icon_default),
             icon_bg=item.get("icon_bg"), title=item.get("title"),
             title_color=item.get("title_color"), title_size=item.get("title_size", 14),
             title_font=item.get("title_font", FONT_HEAD),
             body=item.get("body"), body_color=item.get("body_color"),
             body_size=item.get("body_size", 9))
        cx += int(w) + int(gap)


def stat_row(slide, x, y, w, h, number, label, fill=AZUL, number_color=None, label_color=None):
    """Big-number tile — pass several side by side (see card_row-style
    manual loop, or just call this N times) for a stats row. Figtree Medium
    ~28pt number, light label below."""
    box = rounded_card(slide, x, y, w, h, fill)
    if number_color is None:
        number_color = BLANCO if fill in (AZUL, PANEL_OSCURO) else AZUL_DARK
    if label_color is None:
        label_color = BLANCO if fill in (AZUL, PANEL_OSCURO) else NEGRO
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(int(Inches(0.14)))
    tf.margin_right = Emu(int(Inches(0.14)))
    tf.margin_top = Emu(int(Inches(0.12)))
    tf.margin_bottom = Emu(int(Inches(0.12)))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.size = Pt(28)
    r.font.name = FONT_HEAD_MEDIUM
    r.font.color.rgb = number_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10.5)
    r2.font.name = FONT_CARD_BODY
    r2.font.color.rgb = label_color
    return box


def okr_row(slide, x, y, w, h, tag, title, value, desc, fill=AMARILLO, tag_color=AZUL,
            title_color=GRIS, value_color=NEGRO, desc_color=NEGRO, label_w_frac=0.36):
    """Priority/KR-style row: a short label column (small colored tag +
    gray bold title, e.g. 'KR 1 · FLUJO' / 'Que toda alta nueva se valide')
    to the left of a wide colored stat card (huge value + light description).
    Stack 2-3 of these vertically for an objectives/priorities slide."""
    label_w = Emu(int(int(w) * label_w_frac))
    multi_text(slide, x, y, label_w, h,
               [[{"text": tag, "font": FONT_HEAD, "bold": True, "size": 10.5, "color": tag_color}],
                [{"text": title, "font": FONT_HEAD, "bold": True, "size": 10.5, "color": title_color}]],
               anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=6)
    card_x = Emu(int(x) + int(label_w) + int(Inches(0.15)))
    card_w = Emu(int(x) + int(w) - int(card_x))
    box = rounded_card(slide, card_x, y, card_w, h, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(int(Inches(0.28)))
    tf.margin_right = Emu(int(Inches(0.28)))
    tf.margin_top = tf.margin_bottom = Emu(int(Inches(0.12)))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = FONT_HEAD_MEDIUM
    r.font.color.rgb = value_color
    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.name = FONT_CARD_BODY
    r2.font.color.rgb = desc_color


def team_block(slide, x, y, w_label, h, name, points, fill=AZUL, name_color=None,
               points_color=GRIS, gap=Inches(0.17), points_w=Inches(7.45)):
    """'Quién construye qué' pattern: a narrow colored label card with a team/
    owner name (vertically centered), next to a stacked, unbulleted list of
    ALL-CAPS Figtree Medium capability lines. points: list of str."""
    if name_color is None:
        name_color = BLANCO if fill in (AZUL, PANEL_OSCURO) else NEGRO
    box = rounded_card(slide, x, y, w_label, h, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(int(Inches(0.12)))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.name = FONT_HEAD
    r.font.color.rgb = name_color
    px = Emu(int(x) + int(w_label) + int(gap))
    tb = slide.shapes.add_textbox(px, y, points_w, h)
    ptf = tb.text_frame
    ptf.word_wrap = True
    ptf.margin_left = ptf.margin_right = ptf.margin_top = ptf.margin_bottom = 0
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, pt in enumerate(points):
        pp = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
        pp.space_after = Pt(6)
        rp = pp.add_run()
        rp.text = pt.upper()
        rp.font.size = Pt(11)
        rp.font.name = FONT_HEAD_MEDIUM
        rp.font.color.rgb = points_color


def index_rows(slide, x, y, w, items, row_h=Inches(0.62), number_start=1,
               number_color=GRIS, title_color=NEGRO, body_color=GRIS,
               divider_color=GRIS):
    """Numbered row list with a rule under each — for an agenda/contents
    slide. items: list of (title, body) or (title, None)."""
    cy = int(y)
    for i, item in enumerate(items):
        title_t, body_t = item if isinstance(item, tuple) else (item, None)
        num = f"{number_start + i:02d}"
        textbox(slide, Emu(int(x)), Emu(cy), Inches(0.55), row_h, num, 20, number_color,
                bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        tx = Emu(int(x) + int(Inches(0.68)))
        tw = Emu(int(w) - int(Inches(0.68)))
        if body_t:
            multi_text(slide, tx, Emu(cy), tw, row_h,
                       [[{"text": title_t, "font": FONT_HEAD, "bold": True, "size": 13, "color": title_color}],
                        [{"text": body_t, "font": FONT_CARD_BODY, "bold": False, "size": 10, "color": body_color}]],
                       anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=1)
        else:
            textbox(slide, tx, Emu(cy), tw, row_h, title_t, 13, title_color, bold=True,
                    font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        cy += int(row_h)
        divider(slide, Emu(int(x)), Emu(cy), w, color=divider_color, weight=0.75)
        cy += int(Inches(0.08))
    return cy
